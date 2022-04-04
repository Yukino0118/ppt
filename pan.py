from pptx import Presentation
from pptx.util import Cm 
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
import pandas as pd
from pptx.oxml.xmlchemy import OxmlElement

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def set_cell_border(cell, border_color="000000", border_width='12800'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

#open ppt,create slide
prs = Presentation('1.pptx')
title_only_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_only_slide_layout)
#create object
shapes = slide.shapes

#read excel
df = pd.read_excel('1.xlsx', header = None, index_col = None, engine='openpyxl')

#set table parameter
rows = 12
cols = 10
left = Cm(0)
top = Cm(2)
width = Cm(0.6)
height = Cm(0.1)

#create table
table = shapes.add_table(rows, cols, left, top, width, height).table

#write dataframe to table
for i in range(5):
    table.cell(i,0).text = str(df.at[i,0])
    table.cell(i,0).vertical_anchor = MSO_ANCHOR.MIDDLE



#set table's font
for cell in iter_cells(table):
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(9)
            run.font.name = 'Calibri'
            run.font.color.rgb = RGBColor(255,0,0)

#after setting font, reset table width,height
for i in range(10):
    table.columns[i].width = Cm(1)
    for j in range(12):
        table.rows[j].height = Cm(0.8)

        
#add border(must before drawing)
for i in range(12):
    for j in range(10):
        set_cell_border(table.cell(i,j))

#drawing
table.cell(0,0).fill.solid()
table.cell(0,0).fill.fore_color.rgb = RGBColor(255, 255, 0)






prs.save('1.pptx')